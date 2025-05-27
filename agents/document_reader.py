# agents/universal_document_reader.py

from typing import List, Tuple
from langchain_core.documents import Document
from pptx import Presentation
from docx import Document as DocxDocument
from docx.table import Table
from docx.text.paragraph import Paragraph
import pdfplumber
import pandas as pd
import re
import extract_msg

def read_txt(file) -> Tuple[List[Document], str]:
    text = file.read().decode("utf-8", errors="ignore").strip()
    docs = [Document(page_content=text)] if text else []
    return docs, text

def read_pdf(file) -> Tuple[List[Document], str]:
    with pdfplumber.open(file) as pdf:
        documents = []
        all_text = []
        for i, page in enumerate(pdf.pages):
            text = (page.extract_text() or "").strip()
            if text:
                documents.append(Document(page_content=text, metadata={"page": i+1}))
                all_text.append(text)
    return documents, "\n\n".join(all_text)

def read_docx(file) -> Tuple[List[Document], str]:
    doc = DocxDocument(file)
    full_text = []

    for block in iter_block_items(doc):
        if isinstance(block, str):  # paragraph or list
            full_text.append(block)
        elif isinstance(block, list):  # table
            for row in block:
                full_text.append(" | ".join(cell.strip() for cell in row if cell.strip()))

    text_output = "\n".join(full_text).strip()
    documents = [Document(page_content=text_output)] if text_output else []
    return documents, text_output


def read_pptx(file) -> Tuple[List[Document], str]:
    prs = Presentation(file)
    documents = []
    all_text = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and not shape.has_table:
                slide_text.append(shape.text.strip())

            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    slide_text.append(" | ".join(row_text))  # o "\t".join()

        content = "\n".join(filter(None, slide_text)).strip()
        if content:
            documents.append(Document(page_content=content, metadata={"page": i}))
            all_text.append(content)

    return documents, "\n\n".join(all_text)


def read_msg(file) -> Tuple[List[Document], str]:
    msg = extract_msg.Message(file)

    msg_sender = msg.sender or ""
    msg_to = msg.to or ""
    msg_subject = msg.subject or ""
    msg_date = msg.date or ""
    msg_body = msg.body or ""

    # Opcional: intentar detectar el hilo (respuestas anteriores)
    # Dividimos en partes por separadores comunes en Outlook
    thread_parts = split_email_thread(msg_body)

    # Puedes almacenar cada parte como un subdocumento si quisieras
    full_text = f"""
            [Subject]: {msg_subject}
            [From]: {msg_sender}
            [To]: {msg_to}
            [Date]: {msg_date}

            {msg_body.strip()}
            """.strip()
    
    print("MSG")
    print(full_text)

    doc = Document(
        page_content=full_text,
        metadata={
            "subject": msg_subject,
            "from": msg_sender,
            "to": msg_to,
            "date": msg_date
        }
    )

    return [doc], full_text


def read_xlsx(file) -> Tuple[List[Document], str]:
    xls = pd.ExcelFile(file)
    documents = []
    all_text = []

    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        if df.empty:
            continue
        # Convertimos el DataFrame a texto tabular plano
        text = df.to_csv(sep="\t", index=False)
        content = f"[Sheet: {sheet_name}]\n{text.strip()}"

        documents.append(Document(page_content=content, metadata={"sheet": sheet_name}))
        all_text.append(content)

    return documents, "\n\n".join(all_text)


def read_file_as_documents(file, filename: str) -> Tuple[List[Document], str]:
    ext = filename.lower().split(".")[-1]
    if ext == "pdf":
        return read_pdf(file)
    elif ext == "docx":
        return read_docx(file)
    elif ext == "txt":
        return read_txt(file)
    elif ext == "pptx":
        return read_pptx(file)
    elif ext == "msg":
        return read_msg(file)
    elif ext in ["xls", "xlsx"]:
        return read_xlsx(file)
    else:
        raise ValueError(f"Unsupported file type: {ext}")



def iter_block_items(doc):
    """
    Iterates over paragraphs and tables in the order they appear in the document.
    Returns:
        - str for paragraph
        - list[list[str]] for tables
    """
    from docx.oxml.ns import qn
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl

    parent = doc._element.body
    for child in parent.iterchildren():
        if isinstance(child, CT_P):
            para = Paragraph(child, doc)
            text = para.text.strip()
            if text:
                yield text
        elif isinstance(child, CT_Tbl):
            tbl = Table(child, doc)
            table_data = []
            for row in tbl.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            yield table_data


def split_email_thread(body: str) -> List[str]:
    """
    Intenta dividir un hilo de correo en sus respuestas individuales.
    """
    # Separadores comunes en Outlook y otros clientes
    separators = [
        r"^-+Original Message-+$",
        r"^-+ Forwarded by .+ -+$",
        r"^From:.+\d{4}.*$",
        r"^On .+ wrote:$",
        r"^De: .+\nEnviado el .+\nPara: .+\nAsunto: .+$",  # español
    ]
    pattern = re.compile("|".join(separators), re.MULTILINE | re.IGNORECASE)
    parts = pattern.split(body)
    return [part.strip() for part in parts if part.strip()]